"""
ආනාගේ AI — StudyMate
Rebuilt with custom UI, admin login, Google Search toggle,
separated CSS, all-format support, free API key hints, and more.
"""

import io
import re
import base64
import hashlib
import json
import datetime
from pathlib import Path

import requests
import streamlit as st
import streamlit.components.v1 as components

import PyPDF2
import docx
import google.generativeai as genai

# ── Page config ────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="ආනාගේ AI",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ── Load & inject external CSS ─────────────────────────────────────────────────
CSS_PATH = Path(__file__).parent / "style.css"
if CSS_PATH.exists():
    css_content = CSS_PATH.read_text(encoding="utf-8")
    st.markdown(f"<style>{css_content}</style>", unsafe_allow_html=True)

# ── Secrets ────────────────────────────────────────────────────────────────────
ROOT_FOLDER_ID = st.secrets.get("DRIVE_ROOT_FOLDER_ID", "")
DRIVE_API_KEY  = st.secrets.get("GOOGLE_DRIVE_API_KEY", "")

# Admin credentials — change these in secrets or here
ADMIN_EMAIL    = st.secrets.get("ADMIN_EMAIL", "admin@aana.lk")
ADMIN_PASS_HASH = st.secrets.get(
    "ADMIN_PASS_HASH",
    hashlib.sha256("aana@2025!".encode()).hexdigest()
)

# ── Session state defaults ─────────────────────────────────────────────────────
DEFAULTS = {
    "messages": [],
    "current_subject_id": None,
    "docs": {},
    "subject_name": "",
    "google_search": False,
    "admin_logged_in": False,
    "gemini_key": "",
    "api_key_status": None,       # None | "ok" | "err"
    "free_api_keys": [],          # list of str — managed by admin
    "api_usage_today": {},        # key -> count today
    "usage_date": str(datetime.date.today()),
    "show_admin_modal": False,
    "show_freekeys_modal": False,
    "show_about_modal": False,
    "sidebar_open": False,
    "greeting_shown": False,
}
for k, v in DEFAULTS.items():
    if k not in st.session_state:
        st.session_state[k] = v

# Reset daily usage counter
today = str(datetime.date.today())
if st.session_state.usage_date != today:
    st.session_state.api_usage_today = {}
    st.session_state.usage_date = today

# ── Google Drive helpers ───────────────────────────────────────────────────────
DRIVE_API_BASE = "https://www.googleapis.com/drive/v3/files"
FOLDER_MIME    = "application/vnd.google-apps.folder"

def list_folder(folder_id, api_key):
    folders, files = [], []
    page_token = None
    while True:
        params = {
            "q": f"'{folder_id}' in parents and trashed=false",
            "fields": "nextPageToken,files(id,name,mimeType)",
            "orderBy": "name",
            "pageSize": 100,
            "key": api_key,
        }
        if page_token:
            params["pageToken"] = page_token
        r = requests.get(DRIVE_API_BASE, params=params, timeout=15)
        r.raise_for_status()
        data = r.json()
        if "error" in data:
            raise RuntimeError(data["error"].get("message", str(data["error"])))
        for f in data.get("files", []):
            if f["mimeType"] == FOLDER_MIME:
                folders.append((f["name"], f["id"]))
            else:
                files.append((f["name"], f["id"]))
        page_token = data.get("nextPageToken")
        if not page_token:
            break
    return folders, files

def download_file(file_id, api_key):
    session = requests.Session()
    headers = {"User-Agent": "Mozilla/5.0"}
    url = f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media&key={api_key}"
    r = session.get(url, headers=headers, timeout=40, allow_redirects=True)
    if r.status_code == 403 or len(r.content) < 100:
        url = f"https://drive.google.com/uc?export=download&id={file_id}&confirm=t"
        r = session.get(url, headers=headers, timeout=40, allow_redirects=True)
        if b"Virus scan warning" in r.content[:4000] or "virus scan warning" in r.text[:4000].lower():
            token = re.search(r'confirm=([0-9A-Za-z_\-]+)', r.text)
            uuid  = re.search(r'uuid=([0-9A-Za-z_\-]+)', r.text)
            qs = f"&confirm={token.group(1)}" if token else ""
            qs += f"&uuid={uuid.group(1)}" if uuid else ""
            r = session.get(
                f"https://drive.google.com/uc?export=download&id={file_id}{qs}",
                headers=headers, timeout=40
            )
    r.raise_for_status()
    return r.content

def extract_text(file_bytes, filename):
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".pdf":
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        elif ext in (".docx", ".doc"):
            d = docx.Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in d.paragraphs)
        elif ext in (".pptx",):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts)
            except ImportError:
                return "[PPTX support requires python-pptx]"
        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
                texts = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        texts.append("\t".join(str(c) if c is not None else "" for c in row))
                return "\n".join(texts)
            except ImportError:
                return "[XLSX support requires openpyxl]"
        elif ext in (".txt", ".md", ".csv", ".json"):
            return file_bytes.decode("utf-8", errors="ignore")
        else:
            # Try plain text fallback
            decoded = file_bytes.decode("utf-8", errors="ignore")
            if decoded.strip():
                return decoded
            return f"[Unsupported file format: {ext}]"
    except Exception as e:
        return f"[Could not read {filename}: {e}]"

@st.cache_data(show_spinner=False, ttl=600)
def get_structure(root_id, api_key):
    tree = {}
    years, _ = list_folder(root_id, api_key)
    for year_name, year_id in years:
        tree[year_name] = {}
        sems, _ = list_folder(year_id, api_key)
        for sem_name, sem_id in sems:
            tree[year_name][sem_name] = {}
            subjects, _ = list_folder(sem_id, api_key)
            for subj_name, subj_id in subjects:
                tree[year_name][sem_name][subj_name] = subj_id
    return tree

@st.cache_data(show_spinner=False, ttl=300)
def load_subject_docs(subject_folder_id, api_key):
    SUPPORTED = {".pdf", ".docx", ".doc", ".txt", ".md", ".csv",
                 ".json", ".pptx", ".xlsx", ".xls"}
    _, files = list_folder(subject_folder_id, api_key)
    docs = {}
    for name, fid in files:
        if Path(name).suffix.lower() not in SUPPORTED:
            continue
        try:
            raw  = download_file(fid, api_key)
            text = extract_text(raw, name)
            if text.strip():
                docs[name] = text
        except Exception:
            pass
    return docs

# ── Gemini ─────────────────────────────────────────────────────────────────────
def check_api_key(key):
    try:
        genai.configure(api_key=key)
        m = genai.GenerativeModel("gemini-2.5-flash")
        m.generate_content("Hi")
        return True
    except Exception:
        return False

def ask_gemini(gemini_key, question, docs, use_google=False):
    """Call Gemini; optionally with web search grounding."""
    genai.configure(api_key=gemini_key)

    context = "\n\n".join(
        f"=== Document: {name} ===\n{text[:12000]}"
        for name, text in docs.items()
    )

    note_prompt = f"""You are AI ආනා, a friendly Sri Lankan university study assistant.
Answer using the provided documents. Be thorough and clear.
Format in clear paragraphs. Use bullet points where helpful.
At the very end write: **Sources:** [Doc Name], [Doc Name]
If the answer is not in the documents, say so and suggest enabling Google Search.

---DOCUMENTS---
{context}
---END---

Student question: {question}"""

    google_prompt = f"""You are AI ආනා, a friendly Sri Lankan university study assistant.
First check the provided documents. If the answer isn't fully there, supplement with web knowledge.
Format in clear paragraphs. Cite sources where possible.

---DOCUMENTS---
{context}
---END---

Student question: {question}"""

    prompt = google_prompt if use_google else note_prompt

    if use_google:
        model = genai.GenerativeModel(
            "gemini-2.5-flash",
            tools="google_search_retrieval"
        )
    else:
        model = genai.GenerativeModel("gemini-2.5-flash")

    try:
        resp = model.generate_content(prompt)
        return resp.text
    except Exception:
        # Fallback without grounding
        model = genai.GenerativeModel("gemini-2.5-flash")
        resp = model.generate_content(note_prompt)
        return resp.text

# ── Track API usage ────────────────────────────────────────────────────────────
def track_usage(key):
    masked = key[:8] + "…" if len(key) > 8 else key
    st.session_state.api_usage_today[masked] = (
        st.session_state.api_usage_today.get(masked, 0) + 1
    )

def get_usage(key):
    masked = key[:8] + "…" if len(key) > 8 else key
    return st.session_state.api_usage_today.get(masked, 0)

# ── Render helpers ─────────────────────────────────────────────────────────────
ADMIN_IMG = "assets/admin/admin.jpg"   # served via st.static or base64

def img_b64(path_str):
    p = Path(path_str)
    if p.exists():
        data = p.read_bytes()
        ext  = p.suffix.lstrip(".")
        return f"data:image/{ext};base64,{base64.b64encode(data).decode()}"
    return ""

def render_header():
    gsearch_class = "gsearch-toggle active" if st.session_state.google_search else "gsearch-toggle"
    admin_txt = "🔓 Logout" if st.session_state.admin_logged_in else "🔐 Admin"
    dot_color = "#4ade80" if st.session_state.google_search else "#aaa"

    st.markdown(f"""
<div class="aana-header">
  <div class="hdr-left">
    <button class="hbg-btn" onclick="toggleSidebar()" title="Menu">
      <i class="fa-solid fa-bars"></i>
    </button>
    <button class="hdr-avatar-btn" onclick="openAbout()">
      <img src="{ADMIN_IMG}" alt="ආනා" class="hdr-avatar"
           onerror="this.src='https://ui-avatars.com/api/?name=AI&background=162b56&color=9ecfff&size=38'">
    </button>
    <div class="hdr-title-wrap">
      <a class="hdr-title-link" href="#">
        <div class="hdr-brand">
          <span class="hdr-sinhala">ආනාගේ</span>
          <span class="hdr-ai-text">AI</span>
        </div>
      </a>
      <span class="hdr-sub">Himan Thathuwa Kethala Hiruwa</span>
    </div>
  </div>

  <div class="hdr-right">
    <button class="{gsearch_class}" onclick="toggleGoogle()" title="Toggle Google Search">
      <div class="g-dot" style="background:{dot_color}"></div>
      <span class="gs-label">Google</span>
      <i class="fa-brands fa-google" style="font-size:.7rem"></i>
    </button>
    <button class="admin-btn" onclick="openAdminModal()">
      {admin_txt}
    </button>
  </div>
</div>
""", unsafe_allow_html=True)

def render_footer():
    st.markdown("""
<div class="aana-footer">
  <div class="ftr-copy">© 2025 All Rights Reserved</div>
  <div class="ftr-dev">
    Developed with <span class="ftr-heart">❤️</span>
    <a class="ftr-name" href="https://venurakabojithananda.github.io/" rel="noopener" target="_blank">DSVB</a>
  </div>
</div>
""", unsafe_allow_html=True)

# ── JavaScript bridge ──────────────────────────────────────────────────────────
JS = """
<script>
// ── Sidebar ──
function toggleSidebar() {
  const d = document.getElementById('aana-sidebar');
  const o = document.getElementById('sidebar-overlay');
  if (!d) return;
  const isOpen = d.classList.contains('open');
  d.classList.toggle('open', !isOpen);
  if (o) o.classList.toggle('open', !isOpen);
}
function closeSidebar() {
  const d = document.getElementById('aana-sidebar');
  const o = document.getElementById('sidebar-overlay');
  if (d) d.classList.remove('open');
  if (o) o.classList.remove('open');
}

// ── Tree ──
function toggleNode(id) {
  const ch = document.getElementById('ch-' + id);
  const ic = document.getElementById('ic-' + id);
  if (!ch) return;
  ch.classList.toggle('open');
  if (ic) ic.classList.toggle('open');
}

// ── Google toggle ──
function toggleGoogle() {
  window.parent.postMessage({type:'streamlit:setComponentValue', key:'_google_toggle', value: Date.now()}, '*');
  // Direct DOM update for immediate feedback
  const btn = document.querySelector('.gsearch-toggle');
  if (btn) {
    btn.classList.toggle('active');
    const dot = btn.querySelector('.g-dot');
    if (dot) dot.style.background = btn.classList.contains('active') ? '#4ade80' : '#aaa';
  }
}

// ── Modals ──
function openAbout() {
  document.getElementById('about-modal')?.classList.add('open');
}
function closeAbout() {
  document.getElementById('about-modal')?.classList.remove('open');
}
function openAdminModal() {
  document.getElementById('admin-modal')?.classList.add('open');
}
function closeAdminModal() {
  document.getElementById('admin-modal')?.classList.remove('open');
}
function openFreeKeys() {
  document.getElementById('freekeys-modal')?.classList.add('open');
  closeAdminModal();
}
function closeFreeKeys() {
  document.getElementById('freekeys-modal')?.classList.remove('open');
}

function copyKey(key) {
  navigator.clipboard?.writeText(key).then(() => {
    const btns = document.querySelectorAll('.fk-copy');
    btns.forEach(b => { if (b.dataset.key === key) { b.textContent = '✓ Copied!'; setTimeout(() => b.textContent = 'Copy', 1600); }});
  });
}

// Escape key closes modals
document.addEventListener('keydown', e => {
  if (e.key === 'Escape') {
    closeAbout(); closeAdminModal(); closeFreeKeys(); closeSidebar();
  }
});

// Close on overlay click
document.querySelectorAll('.modal-overlay').forEach(el => {
  el.addEventListener('click', function(e) { if (e.target === this) this.classList.remove('open'); });
});
</script>
"""

# ═══════════════════════════════════════════════════════════════════════════════
#   MAIN APP
# ═══════════════════════════════════════════════════════════════════════════════

# ── Handle URL query actions (Google toggle via Streamlit hidden buttons) ──────
if "google_search_toggle" in st.query_params:
    st.session_state.google_search = not st.session_state.google_search
    st.query_params.clear()

# ── Inject JS ──────────────────────────────────────────────────────────────────
st.markdown(JS, unsafe_allow_html=True)

# ── Header ─────────────────────────────────────────────────────────────────────
render_header()

# ── App layout ─────────────────────────────────────────────────────────────────
if not DRIVE_API_KEY or not ROOT_FOLDER_ID:
    st.error("⚠️ DRIVE_ROOT_FOLDER_ID and GOOGLE_DRIVE_API_KEY must be set in Streamlit secrets.")
    st.stop()

# Load Drive structure
try:
    structure = get_structure(ROOT_FOLDER_ID, DRIVE_API_KEY)
except Exception as e:
    st.error(f"Drive API error: {e}")
    st.stop()

# ── Sidebar overlay ────────────────────────────────────────────────────────────
st.markdown('<div class="sidebar-overlay" id="sidebar-overlay" onclick="closeSidebar()"></div>', unsafe_allow_html=True)

# ── Build sidebar HTML ─────────────────────────────────────────────────────────
def render_sidebar(structure):
    # API key section
    cur_key = st.session_state.gemini_key
    status_html = ""
    if st.session_state.api_key_status == "ok":
        status_html = '<div class="api-status ok"><i class="fa-solid fa-circle-check"></i> API key valid ✓</div>'
    elif st.session_state.api_key_status == "err":
        status_html = '<div class="api-status err"><i class="fa-solid fa-circle-xmark"></i> Invalid key — try another</div>'

    # Build folder tree — default expand Year 01 only
    tree_html = ""
    for yi, (year_name, sems) in enumerate(structure.items()):
        y_open = yi == 0
        y_open_cls = "open" if y_open else ""
        tree_html += f"""
        <div class="tree-node">
          <button class="tree-item" onclick="toggleNode('y{yi}')">
            <span class="tree-chevron {'open' if y_open else ''}" id="ic-y{yi}"><i class="fa-solid fa-chevron-right"></i></span>
            <i class="fa-solid fa-graduation-cap" style="color:#58a6ff;font-size:.8rem"></i>
            {year_name}
          </button>
          <div class="tree-children {y_open_cls}" id="ch-y{yi}">
        """
        for si, (sem_name, subjects) in enumerate(sems.items()):
            s_open_cls = ""  # semesters collapsed by default
            tree_html += f"""
            <div class="tree-node">
              <button class="tree-item" onclick="toggleNode('y{yi}s{si}')">
                <span class="tree-chevron" id="ic-y{yi}s{si}"><i class="fa-solid fa-chevron-right"></i></span>
                <i class="fa-solid fa-book-open" style="color:#ffd84d;font-size:.75rem"></i>
                {sem_name}
              </button>
              <div class="tree-children" id="ch-y{yi}s{si}">
            """
            for subj_name, subj_id in subjects.items():
                sel_cls = "selected" if subj_id == st.session_state.current_subject_id else ""
                tree_html += f"""
                <button class="tree-leaf {sel_cls}"
                        onclick="selectSubject('{subj_id}', '{subj_name.replace("'", "\\'")}')">
                  <i class="fa-solid fa-file-lines" style="font-size:.72rem"></i>
                  {subj_name}
                </button>
                """
            tree_html += "</div></div>"
        tree_html += "</div></div>"

    # Free API keys hint btn
    free_keys_btn = """
    <button class="free-keys-btn" onclick="openFreeKeys()">
      <span class="fk-icon">🔑</span>
      <span>තාමත් API Key එකක් නැද්ද..?</span>
    </button>
    """

    sidebar_html = f"""
    <div class="sidebar-drawer" id="aana-sidebar">
      <div class="sidebar-hdr">
        <span class="sidebar-title">ආනාගේ AI</span>
        <button class="sidebar-close" onclick="closeSidebar()"><i class="fa-solid fa-xmark"></i></button>
      </div>
      <div class="sidebar-body">

        {free_keys_btn}

        <div class="api-section">
          <div class="api-section-title"><i class="fa-solid fa-key"></i> Gemini API Key</div>
          <div class="api-input-wrap">
            <input class="api-input" id="api-key-input" type="password"
                   placeholder="AIza..." value="{cur_key}"
                   oninput="apiKeyChanged(this.value)">
          </div>
          {status_html}
          <button class="api-check-btn" onclick="checkApiKey()">
            <i class="fa-solid fa-circle-check"></i> Check Key
          </button>
          <div style="margin-top:.45rem;font-size:.7rem;color:#484f58">
            Get free key →
            <a href="https://aistudio.google.com/app/apikey" target="_blank"
               style="color:#58a6ff">aistudio.google.com</a>
          </div>
        </div>

        <div class="tree-section">
          <div class="tree-section-title"><i class="fa-solid fa-folder-open"></i> Your Subject</div>
          {tree_html}
        </div>

        <div id="selected-subject-info" style="display:none;background:#111827;border:1px solid #1f2d42;border-radius:10px;padding:.65rem .85rem;margin-bottom:.5rem;">
          <div style="font-size:.7rem;color:#58a6ff;font-weight:800;text-transform:uppercase;letter-spacing:.06em;margin-bottom:.3rem;">Selected</div>
          <div id="selected-subject-name" style="color:#f0f6fc;font-size:.85rem;font-weight:700;"></div>
        </div>

        <button class="sidebar-load-btn" id="load-btn" onclick="loadSubject()">
          <i class="fa-solid fa-download"></i> Load Subject Docs
        </button>
        <button class="sidebar-clear-btn" onclick="newChat()">
          <i class="fa-solid fa-rotate-left"></i> Clear Chat
        </button>

        <div style="margin-top:.75rem;font-size:.68rem;color:#484f58;text-align:center;line-height:1.6;">
          Docs refresh every 5 min.<br>Drop new files in Drive anytime.
        </div>
      </div>
    </div>
    """
    st.markdown(sidebar_html, unsafe_allow_html=True)

render_sidebar(structure)

# ── Sidebar JS (subject selection, API key check) ──────────────────────────────
st.markdown("""
<script>
let _selectedSubjectId = null;
let _selectedSubjectName = null;

function selectSubject(id, name) {
  _selectedSubjectId = id;
  _selectedSubjectName = name;
  // Highlight
  document.querySelectorAll('.tree-leaf').forEach(el => el.classList.remove('selected'));
  event.currentTarget.classList.add('selected');
  // Show info
  const info = document.getElementById('selected-subject-info');
  const nm   = document.getElementById('selected-subject-name');
  if (info) info.style.display = '';
  if (nm)   nm.textContent = name;
}

function apiKeyChanged(val) {
  // Store in hidden input that Streamlit reads
  const h = document.getElementById('_api_key_hidden');
  if (h) h.value = val;
}

function checkApiKey() {
  const val = document.getElementById('api-key-input')?.value || '';
  setParam('check_key', val);
}

function loadSubject() {
  if (!_selectedSubjectId) { alert('Please select a subject first!'); return; }
  setParam('load_subj', _selectedSubjectId + '|' + _selectedSubjectName);
}

function newChat() {
  setParam('new_chat', '1');
}

function setParam(key, val) {
  const url = new URL(window.location.href);
  url.searchParams.set(key, val);
  window.location.href = url.toString();
}

// Google toggle
function toggleGoogle() {
  const url = new URL(window.location.href);
  url.searchParams.set('google_toggle', '1');
  window.location.href = url.toString();
}
</script>
""", unsafe_allow_html=True)

# ── Handle query params (actions) ──────────────────────────────────────────────
params = st.query_params

if "google_toggle" in params:
    st.session_state.google_search = not st.session_state.google_search
    st.query_params.clear()
    st.rerun()

if "new_chat" in params:
    st.session_state.messages = []
    st.session_state.greeting_shown = False
    st.query_params.clear()
    st.rerun()

if "check_key" in params:
    key_to_check = params["check_key"]
    st.session_state.gemini_key = key_to_check
    if key_to_check:
        with st.spinner("Checking API key…"):
            ok = check_api_key(key_to_check)
        st.session_state.api_key_status = "ok" if ok else "err"
    st.query_params.clear()
    st.rerun()

if "load_subj" in params:
    raw = params["load_subj"]
    parts = raw.split("|", 1)
    subj_id   = parts[0]
    subj_name = parts[1] if len(parts) > 1 else "Subject"
    if subj_id != st.session_state.current_subject_id:
        st.session_state.messages = []
        st.session_state.greeting_shown = False
        st.session_state.current_subject_id = subj_id
        st.session_state.subject_name = subj_name
        st.session_state.docs = {}
    with st.spinner(f"Loading docs for {subj_name}…"):
        st.session_state.docs = load_subject_docs(subj_id, DRIVE_API_KEY)
    st.query_params.clear()
    st.rerun()

# ── Main chat area ─────────────────────────────────────────────────────────────
docs = st.session_state.docs
subj_name = st.session_state.subject_name

st.markdown('<div class="app-layout">', unsafe_allow_html=True)
st.markdown('<div class="main-wrap">', unsafe_allow_html=True)

if not docs:
    # Hero / landing
    if not st.session_state.current_subject_id:
        st.markdown(f"""
        <div class="hero-wrap">
          <div class="hero-avatar-wrap">
            <img src="{ADMIN_IMG}" alt="AI ආනා"
                 onerror="this.src='https://ui-avatars.com/api/?name=AI&background=162b56&color=9ecfff&size=90'">
          </div>
          <div class="hero-title">ආනාගේ <span>AI</span></div>
          <div class="hero-sub">ඔබේ lecture notes ගැන ඕනෙ ප්‍රශ්නයක් අහන්නකෝ..! 🎓</div>
          <div class="hero-steps">
            <div class="hero-step"><span class="hero-step-num">1</span> API Key paste කරන්න</div>
            <div class="hero-step"><span class="hero-step-num">2</span> Subject select කරන්න</div>
            <div class="hero-step"><span class="hero-step-num">3</span> Load කරලා අහන්න</div>
          </div>
        </div>
        """, unsafe_allow_html=True)
    else:
        # Subject selected but not loaded
        st.markdown(f"""
        <div class="load-subject-banner">
          <div class="lsb-icon">📂</div>
          <h3>Load Subject Docs</h3>
          <p>"{subj_name}" select කරලා තියෙනවා.<br>
          Sidebar එකේ <strong>Load Subject Docs</strong> button එක press කරන්නකෝ!</p>
        </div>
        """, unsafe_allow_html=True)

else:
    # Status bar
    doc_names = list(docs.keys())
    tags = "".join(f"<span class='doc-tag' title='{n}'>{n}</span>" for n in doc_names[:4])
    if len(doc_names) > 4:
        tags += f"<span class='doc-tag'>+{len(doc_names)-4} more</span>"
    gsearch_badge = (
        '<span style="color:#4ade80;font-size:.75rem;font-weight:700">🌐 Google ON</span>'
        if st.session_state.google_search else ""
    )
    st.markdown(f"""
    <div class="status-bar">
      <span>📂</span>
      <span class="status-subject">{subj_name}</span>
      <span style="color:var(--border)">|</span>
      <span>{len(doc_names)} doc{'s' if len(doc_names)!=1 else ''}</span>
      <span style="color:var(--border)">|</span>
      {tags}
      {gsearch_badge}
    </div>
    """, unsafe_allow_html=True)

    # Greeting
    if not st.session_state.greeting_shown:
        st.session_state.messages = []
        greeting = (
            "හෙලෝ සුද්දා, කෝමද, සැපේද ඉන්නේ..? 😊\n\n"
            f"You can now ask any question regarding the uploaded notes of **{subj_name}**.. "
            "Can't find your answer? Simply enable Google Search 🌐"
        )
        st.session_state.messages.append({"role": "assistant", "content": greeting, "is_greeting": True})
        st.session_state.greeting_shown = True

    # Render messages
    st.markdown('<div class="chat-area" id="chat-area">', unsafe_allow_html=True)
    for msg in st.session_state.messages:
        if msg["role"] == "user":
            st.markdown(f"""
            <div class="msg-user">
              <div class="msg-user-bubble">{msg['content']}</div>
            </div>""", unsafe_allow_html=True)
        else:
            content_html = msg['content'].replace('\n', '<br>')
            is_greeting  = msg.get("is_greeting", False)
            bubble_cls   = "greeting-msg" if is_greeting else "msg-ai-bubble"
            avatar_html  = f"""
              <div class="msg-ai-avatar">
                <img src="{ADMIN_IMG}" alt="AI ආනා"
                     onerror="this.style.opacity=0">
              </div>"""

            # Check if answer suggests enabling Google
            show_google_chip = (
                not st.session_state.google_search
                and not is_greeting
                and any(x in msg['content'].lower() for x in [
                    "not in the documents", "not found", "cannot find",
                    "no information", "document doesn't", "isn't in"
                ])
            )

            google_chip = ""
            if show_google_chip:
                google_chip = """
                <br>
                <button class="google-suggest" onclick="toggleGoogle()">
                  <i class="fa-brands fa-google"></i> Enable Google Search for a better answer
                </button>"""

            st.markdown(f"""
            <div class="msg-ai">
              {avatar_html}
              <div>
                <div class="msg-ai-name">AI ආනා</div>
                <div class="{bubble_cls}">{content_html}{google_chip}</div>
              </div>
            </div>""", unsafe_allow_html=True)

    st.markdown('</div>', unsafe_allow_html=True)

    # Auto-scroll
    st.markdown("""
    <script>
    setTimeout(() => {
      const ca = document.getElementById('chat-area');
      if (ca) ca.scrollTop = ca.scrollHeight;
    }, 200);
    </script>
    """, unsafe_allow_html=True)

    # ── Chat input ──
    col1, col2, col3 = st.columns([1, 10, 1])
    with col2:
        question = st.chat_input(f"Ask about {subj_name}…")

    if question:
        key = st.session_state.gemini_key
        if not key:
            st.error("⚠️ Sidebar-এ API key paste කරන්නකෝ first!")
        else:
            st.session_state.messages.append({"role": "user", "content": question})
            with st.spinner("AI ආනා thinking… 🤔"):
                try:
                    if st.session_state.google_search:
                        st.info("🌐 Google Search enabled — may enhance your answer", icon="ℹ️")
                    answer = ask_gemini(key, question, docs, st.session_state.google_search)
                    track_usage(key)
                    st.session_state.messages.append({"role": "assistant", "content": answer})
                    st.rerun()
                except Exception as e:
                    err = str(e)
                    if "API_KEY_INVALID" in err or "invalid" in err.lower():
                        st.error("❌ Invalid Gemini API key. Get a free one at aistudio.google.com")
                        st.session_state.api_key_status = "err"
                    else:
                        st.error(f"Gemini error: {err}")

st.markdown('</div></div>', unsafe_allow_html=True)  # close main-wrap + app-layout

# ══════════════════════════════════════
#   MODALS
# ══════════════════════════════════════

# ── About / profile modal ──────────────────────────────────────────────────────
st.markdown(f"""
<div class="modal-overlay" id="about-modal" onclick="if(event.target===this)closeAbout()">
  <div class="modal-box">
    <button class="modal-close" onclick="closeAbout()"><i class="fa-solid fa-xmark"></i></button>
    <div class="modal-avatar-wrap" style="cursor:pointer">
      <img src="{ADMIN_IMG}" class="modal-avatar" alt="ආනා"
           onerror="this.style.opacity=0">
    </div>
    <div class="modal-name">ආනා</div>
    <div class="modal-role">B.Ed Hons English | Dip. English Lang &amp; Lit.</div>
    <a class="wa-btn" href="https://wa.me/94784892024?text=ආනා%20අයියේ..%20පොඩි%20සීන්%20එකක්%20බං%20support%20එකක්%20දෙන්නකෝ.."
       rel="noopener" target="_blank">
      <i class="fa-brands fa-whatsapp"></i> Contact ආනා
    </a>
  </div>
</div>
""", unsafe_allow_html=True)

# ── Admin modal ────────────────────────────────────────────────────────────────
if not st.session_state.admin_logged_in:
    st.markdown("""
<div class="modal-overlay" id="admin-modal" onclick="if(event.target===this)closeAdminModal()">
  <div class="modal-box">
    <button class="modal-close" onclick="closeAdminModal()"><i class="fa-solid fa-xmark"></i></button>
    <div class="modal-title">🔐 Admin Login</div>
    <div class="modal-sub">Admin-only area</div>
    <form onsubmit="doAdminLogin(event)">
      <input class="admin-field" id="admin-email" type="email" placeholder="Email">
      <input class="admin-field" id="admin-pass"  type="password" placeholder="Password">
      <div class="admin-err" id="admin-err" style="display:none">❌ Invalid credentials</div>
      <button type="submit" class="admin-login-btn">Login</button>
    </form>
  </div>
</div>
""", unsafe_allow_html=True)

    # Admin login JS — posts to Streamlit via query param
    admin_hash = ADMIN_PASS_HASH
    admin_email = ADMIN_EMAIL
    st.markdown(f"""
<script>
function doAdminLogin(e) {{
  e.preventDefault();
  const email = document.getElementById('admin-email').value.trim();
  const pass  = document.getElementById('admin-pass').value;
  // Simple hash check (SHA-256 via SubtleCrypto)
  crypto.subtle.digest('SHA-256', new TextEncoder().encode(pass)).then(buf => {{
    const hex = Array.from(new Uint8Array(buf)).map(b=>b.toString(16).padStart(2,'0')).join('');
    if (email === '{admin_email}' && hex === '{admin_hash}') {{
      window.location.href = window.location.pathname + '?admin_login=1';
    }} else {{
      const err = document.getElementById('admin-err');
      if (err) {{ err.style.display=''; setTimeout(()=>err.style.display='none', 3000); }}
    }}
  }});
}}
</script>
""", unsafe_allow_html=True)

    if "admin_login" in params:
        st.session_state.admin_logged_in = True
        st.query_params.clear()
        st.rerun()

else:
    # Admin logged in — show panel
    free_keys = st.session_state.free_api_keys
    keys_html = ""
    for i, k in enumerate(free_keys):
        usage = get_usage(k)
        keys_html += f"""
        <div class="admin-api-item">
          <input class="admin-api-input" value="{k}" onchange="updateKey({i}, this.value)">
          <span style="font-size:.68rem;color:#484f58;white-space:nowrap">{usage}x today</span>
          <button class="admin-api-del" onclick="deleteKey({i})" title="Remove">
            <i class="fa-solid fa-trash"></i>
          </button>
        </div>"""

    st.markdown(f"""
<div class="modal-overlay" id="admin-modal" onclick="if(event.target===this)closeAdminModal()">
  <div class="admin-panel">
    <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:.75rem">
      <div class="admin-panel-title">⚙️ Admin Panel</div>
      <button class="modal-close" style="position:static" onclick="closeAdminModal()">
        <i class="fa-solid fa-xmark"></i>
      </button>
    </div>
    <div style="font-size:.72rem;color:#58a6ff;font-weight:800;text-transform:uppercase;
                letter-spacing:.06em;margin-bottom:.5rem">Free API Keys</div>
    {keys_html}
    <button class="admin-add-btn" onclick="addKey()">
      <i class="fa-solid fa-plus"></i> Add API Key
    </button>
    <button class="admin-logout-btn" onclick="adminLogout()">
      <i class="fa-solid fa-right-from-bracket"></i> Logout
    </button>
  </div>
</div>
""", unsafe_allow_html=True)

    # Admin actions via query params
    if "admin_logout" in params:
        st.session_state.admin_logged_in = False
        st.query_params.clear()
        st.rerun()
    if "admin_add_key" in params:
        new_key = params["admin_add_key"]
        if new_key and new_key not in st.session_state.free_api_keys:
            st.session_state.free_api_keys.append(new_key)
        st.query_params.clear()
        st.rerun()
    if "admin_del_key" in params:
        idx = int(params["admin_del_key"])
        if 0 <= idx < len(st.session_state.free_api_keys):
            st.session_state.free_api_keys.pop(idx)
        st.query_params.clear()
        st.rerun()
    if "admin_upd_key" in params:
        raw = params["admin_upd_key"]
        parts = raw.split("|", 1)
        if len(parts) == 2:
            idx, val = int(parts[0]), parts[1]
            if 0 <= idx < len(st.session_state.free_api_keys):
                st.session_state.free_api_keys[idx] = val
        st.query_params.clear()
        st.rerun()

    st.markdown("""
<script>
function addKey() {
  const k = prompt('Paste new API key:');
  if (k && k.trim()) setParam('admin_add_key', k.trim());
}
function deleteKey(i) {
  if (confirm('Remove this key?')) setParam('admin_del_key', i);
}
function updateKey(i, val) {
  setParam('admin_upd_key', i + '|' + val);
}
function adminLogout() { setParam('admin_logout', '1'); }
</script>
""", unsafe_allow_html=True)

# ── Free Keys Modal ────────────────────────────────────────────────────────────
free_keys = st.session_state.free_api_keys
if not free_keys:
    # Show placeholder if no admin keys yet
    keys_content = '<div style="color:#484f58;font-size:.8rem;text-align:center;padding:.5rem 0">No free keys added yet. Admin can add them.</div>'
else:
    keys_content = ""
    for k in free_keys:
        usage = get_usage(k)
        masked = k[:12] + "…" + k[-4:] if len(k) > 18 else k
        keys_content += f"""
        <div class="fk-item">
          <span class="fk-key">{masked}</span>
          <button class="fk-copy" data-key="{k}" onclick="copyKey('{k}')">Copy</button>
        </div>
        <div class="fk-usage">Used {usage}x today</div>
        """

st.markdown(f"""
<div class="modal-overlay" id="freekeys-modal" onclick="if(event.target===this)closeFreeKeys()">
  <div class="fk-modal-box">
    <button class="modal-close" style="background:rgba(255,255,255,.12);color:#f0f6fc"
            onclick="closeFreeKeys()"><i class="fa-solid fa-xmark"></i></button>
    <div class="fk-title">🔑 Free API Keys</div>
    <div class="fk-warn">
      මෙවුවා හැමෝම use කරනවා limit වැදිලා ඇති සමහරවිට 😠<br>
      Video එක බලලා තමන්ටම කියලා එකක් හදාගනිං API හිඟන්නා 😤
    </div>
    {keys_content}
    <a class="fk-video-link" href="https://youtu.be/YOUR_VIDEO_LINK_HERE" target="_blank" rel="noopener">
      <i class="fa-brands fa-youtube"></i> How to get your own free API key →
    </a>
  </div>
</div>
""", unsafe_allow_html=True)

# ── Footer ─────────────────────────────────────────────────────────────────────
render_footer()
